"""
小棲蘭森林農場 V2 網站架構說明書 PPTX 生成腳本
執行方式：python generate_pptx.py
需要：pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ── 色彩定義 ──────────────────────────────────────────
COLOR_FOREST_GREEN  = RGBColor(0x06, 0x4E, 0x3B)   # emerald-900
COLOR_MID_GREEN     = RGBColor(0x05, 0x96, 0x69)   # emerald-700
COLOR_LIGHT_GREEN   = RGBColor(0xD1, 0xFA, 0xE5)   # emerald-100
COLOR_BEIGE         = RGBColor(0xFA, 0xFA, 0xF9)   # stone-50
COLOR_CHARCOAL      = RGBColor(0x1F, 0x29, 0x37)   # gray-800
COLOR_RED           = RGBColor(0xB9, 0x1C, 0x1C)   # red-700
COLOR_WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY          = RGBColor(0x6B, 0x72, 0x80)   # gray-500
COLOR_AMBER         = RGBColor(0xD9, 0x77, 0x06)   # amber-600
COLOR_DARK_BG       = RGBColor(0x02, 0x2C, 0x22)   # very dark green

# ── 投影片尺寸（16:9 寬螢幕）──────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color: RGBColor):
    """設定投影片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=COLOR_CHARCOAL,
                 align=PP_ALIGN.LEFT, font_name="微軟正黑體"):
    """新增文字方塊"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    """新增矩形色塊"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_bullet_text(slide, items, left, top, width, height,
                    font_size=14, color=COLOR_CHARCOAL, bullet="▸ "):
    """新增項目符號文字方塊"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "微軟正黑體"
    return txBox


# ════════════════════════════════════════════════════════
# 建立簡報
# ════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # 完全空白版面


# ════════════════════════════════════════════════════════
# 投影片 1：封面
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, COLOR_DARK_BG)

# 左側深綠色塊
add_rect(slide, Inches(0), Inches(0), Inches(5), SLIDE_H, COLOR_FOREST_GREEN)

# 裝飾線
add_rect(slide, Inches(5), Inches(2.8), Inches(0.06), Inches(1.8), COLOR_MID_GREEN)

# 主標題
add_text_box(slide, "🌲 小棲蘭森林農場",
             Inches(0.4), Inches(1.2), Inches(4.2), Inches(0.8),
             font_size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "V2 靜態 MVP 網站",
             Inches(0.4), Inches(2.0), Inches(4.2), Inches(0.7),
             font_size=28, bold=True, color=COLOR_LIGHT_GREEN, align=PP_ALIGN.CENTER)

add_text_box(slide, "架構規劃說明書",
             Inches(0.4), Inches(2.7), Inches(4.2), Inches(0.7),
             font_size=28, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# 副標題（右側）
add_text_box(slide, "CWW 架構對齊 · 模組化設計 · 100% Mobile-First",
             Inches(5.4), Inches(2.5), Inches(7.5), Inches(0.6),
             font_size=16, bold=False, color=COLOR_LIGHT_GREEN, align=PP_ALIGN.LEFT)

add_text_box(slide, "Cloudways + WordPress + WooCommerce\n無縫遷移路徑規劃",
             Inches(5.4), Inches(3.2), Inches(7.5), Inches(1.0),
             font_size=20, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)

# 技術棧標籤
tags = ["HTML5", "Tailwind CSS", "FontAwesome 6", "Noto Serif TC", "Vanilla ES6"]
for i, tag in enumerate(tags):
    add_rect(slide, Inches(5.4 + i * 1.5), Inches(4.5), Inches(1.35), Inches(0.4), COLOR_MID_GREEN)
    add_text_box(slide, tag,
                 Inches(5.4 + i * 1.5), Inches(4.5), Inches(1.35), Inches(0.4),
                 font_size=10, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# 日期
add_text_box(slide, "2026.06.23  ·  建安農業有限公司",
             Inches(5.4), Inches(6.5), Inches(7.5), Inches(0.5),
             font_size=12, color=COLOR_GRAY, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════
# 投影片 2：專案概覽 & 核心工程指令
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, COLOR_BEIGE)

# 頂部標題列
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), COLOR_FOREST_GREEN)
add_text_box(slide, "01  專案概覽 & 核心工程指令（Architecture Law）",
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=COLOR_WHITE)

# 左欄：專案目標
add_rect(slide, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.8), COLOR_WHITE,
         line_color=COLOR_LIGHT_GREEN, line_width=Pt(1))
add_text_box(slide, "🎯  專案目標",
             Inches(0.5), Inches(1.45), Inches(5.4), Inches(0.5),
             font_size=16, bold=True, color=COLOR_FOREST_GREEN)
add_bullet_text(slide, [
    "為「小棲蘭森林農場」打造 5 頁靜態 MVP 原型",
    "100% Mobile-First RWD，桌機與手機皆完美呈現",
    "模組化元件設計，非技術農場主人可輕鬆維護",
    "對齊 CWW（Cloudways + WordPress + WooCommerce）",
    "所有連結使用扁平路徑，直接對應 WP URL Slug",
    "純 Vanilla ES6 JS，無框架依賴，解耦設計",
], Inches(0.5), Inches(2.0), Inches(5.4), Inches(4.5), font_size=13)

# 右欄：三大架構法則
add_rect(slide, Inches(6.5), Inches(1.3), Inches(6.5), Inches(5.8), COLOR_WHITE,
         line_color=COLOR_LIGHT_GREEN, line_width=Pt(1))
add_text_box(slide, "⚖️  三大架構法則（Architecture Law）",
             Inches(6.7), Inches(1.45), Inches(6.1), Inches(0.5),
             font_size=16, bold=True, color=COLOR_FOREST_GREEN)

laws = [
    ("① CWW 對齊", "每個頁面 HTML 結構高度模組化（元件化），媒體與商品容器完全解耦，農場主人可在 WP Gutenberg 後台直接替換文字、圖片、YouTube 連結，如同使用 Facebook 般直覺。"),
    ("② 乾淨動態路由", "所有導覽錨點連結使用扁平、一致的路徑（events.html / shop.html），可直接對應 WordPress URL Slug，零修改成本遷移。"),
    ("③ Mobile-First RWD", "每個元素皆以手機優先設計，使用 Tailwind CSS 響應式前綴（md: / lg:）確保桌機與手機皆有頂級視覺體驗。"),
]
y_pos = 2.05
for title, desc in laws:
    add_rect(slide, Inches(6.7), Inches(y_pos), Inches(6.0), Inches(0.35), COLOR_LIGHT_GREEN)
    add_text_box(slide, title,
                 Inches(6.8), Inches(y_pos), Inches(5.8), Inches(0.35),
                 font_size=13, bold=True, color=COLOR_FOREST_GREEN)
    add_text_box(slide, desc,
                 Inches(6.7), Inches(y_pos + 0.38), Inches(6.0), Inches(0.9),
                 font_size=11, color=COLOR_CHARCOAL)
    y_pos += 1.45


# ════════════════════════════════════════════════════════
# 投影片 3：5 頁面檔案架構
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, COLOR_BEIGE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), COLOR_FOREST_GREEN)
add_text_box(slide, "02  5 頁面檔案架構（File Architecture）",
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=COLOR_WHITE)

pages = [
    ("index.html", "全螢幕歡迎頁", COLOR_FOREST_GREEN,
     ["全滿版背景圖（w-full h-screen overflow-hidden）",
      "整個 <body> 包在 <a> 內，點擊任意處跳轉",
      "右下角品牌簽名：毛筆字「小棲蘭」+ 紅色印章「建安」",
      "雙重安全邊距防止窄螢幕裁切"]),
    ("events.html", "最新活動頁（主 Hub）", COLOR_MID_GREEN,
     ["全域黃金導覽列（桌機 + 手機 RWD）",
      "Hero Carousel 輪播（純 Tailwind + Vanilla ES6）",
      "活動卡片 Grid（md:grid-cols-2）",
      "報名按鈕連結 Google Forms（target=_blank）",
      "注意事項 & 交通資訊區塊"]),
    ("shop.html", "森林農產品頁（B2C 電商）", RGBColor(0x92, 0x40, 0x0E),
     ["100% B2C 專屬，移除所有 B2B 詢價表單",
      "3 件核心商品卡片（高信任徽章：貨到付款 / 免運）",
      "加入購物車 Badge 計數器（Vanilla ES6 解耦）",
      "支付閘道預留 Hook（Visa / Mastercard disabled）"]),
    ("blog.html", "小棲蘭日誌（SEO 引擎）", RGBColor(0x0E, 0x4C, 0x6B),
     ["精選置頂文章 + 3 欄 Timeline Grid",
      "防爆版 YouTube 嵌入（aspect-video 嚴格寬高比）",
      "料理 / 養生文章末尾情境式 CTA → shop.html",
      "電子報訂閱留存機制"]),
    ("cart.html", "購物車結帳頁（轉換漏斗）", RGBColor(0x4C, 0x0E, 0x0E),
     ["3 欄 RWD 佈局（商品明細 col-span-2 + 訂單摘要）",
      "FontAwesome 圓形 +/- 數量控制器",
      "貨到付款預設強制預選（disabled 線上刷卡）",
      "信任徽章：有機認證 / 24h 出貨 / 退款保障",
      "結帳成功 Modal 彈窗"]),
]

col_w = Inches(2.4)
for i, (filename, title, color, bullets) in enumerate(pages):
    x = Inches(0.3 + i * 2.6)
    add_rect(slide, x, Inches(1.2), col_w, Inches(0.55), color)
    add_text_box(slide, filename,
                 x, Inches(1.22), col_w, Inches(0.3),
                 font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 x, Inches(1.52), col_w, Inches(0.3),
                 font_size=10, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, x, Inches(1.75), col_w, Inches(5.4), COLOR_WHITE,
             line_color=color, line_width=Pt(1.5))
    add_bullet_text(slide, bullets,
                    Inches(0.35 + i * 2.6), Inches(1.85), Inches(2.3), Inches(5.0),
                    font_size=10, color=COLOR_CHARCOAL, bullet="• ")


# ════════════════════════════════════════════════════════
# 投影片 4：設計系統 & 色彩規範
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, COLOR_BEIGE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), COLOR_FOREST_GREEN)
add_text_box(slide, "03  設計系統 & 色彩規範（Design System）",
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=COLOR_WHITE)

# 色彩展示
colors_info = [
    ("深森林綠", "#064E3B", COLOR_FOREST_GREEN, "主色 · 導覽列 · 標題"),
    ("中森林綠", "#059669", COLOR_MID_GREEN, "按鈕 · CTA · 強調"),
    ("大地米白", "#FAFAF9", RGBColor(0xFA, 0xFA, 0xF9), "頁面背景 · 卡片底色"),
    ("警示紅", "#B91C1C", COLOR_RED, "價格 · 促銷 · 印章"),
    ("炭灰色", "#1F2937", COLOR_CHARCOAL, "正文 · 說明文字"),
    ("琥珀橙", "#D97706", COLOR_AMBER, "標籤 · 免運橫幅"),
]

add_text_box(slide, "🎨  色彩規範",
             Inches(0.4), Inches(1.2), Inches(4), Inches(0.4),
             font_size=16, bold=True, color=COLOR_FOREST_GREEN)

for i, (name, hex_val, rgb, usage) in enumerate(colors_info):
    y = Inches(1.7 + i * 0.85)
    add_rect(slide, Inches(0.4), y, Inches(1.2), Inches(0.6), rgb,
             line_color=RGBColor(0xD1, 0xD5, 0xDB), line_width=Pt(0.5))
    add_text_box(slide, name,
                 Inches(1.7), y, Inches(1.5), Inches(0.3),
                 font_size=13, bold=True, color=COLOR_CHARCOAL)
    add_text_box(slide, f"{hex_val}  ·  {usage}",
                 Inches(1.7), Inches(1.7 + i * 0.85 + 0.3), Inches(4.5), Inches(0.3),
                 font_size=11, color=COLOR_GRAY)

# 字型規範
add_text_box(slide, "🔤  字型規範",
             Inches(6.5), Inches(1.2), Inches(6.5), Inches(0.4),
             font_size=16, bold=True, color=COLOR_FOREST_GREEN)

font_info = [
    ("主要字型", "Noto Serif TC", "Google Fonts CDN", "標題 / 品牌文字 / 強調"),
    ("系統字型", "微軟正黑體 / PingFang TC", "系統預設", "正文 / 說明 / 按鈕"),
    ("圖示字型", "FontAwesome 6 Free", "cdnjs CDN", "所有 UI 圖示"),
]

y_pos = 1.7
for category, font, source, usage in font_info:
    add_rect(slide, Inches(6.5), Inches(y_pos), Inches(6.5), Inches(1.1), COLOR_WHITE,
             line_color=COLOR_LIGHT_GREEN, line_width=Pt(1))
    add_text_box(slide, category,
                 Inches(6.7), Inches(y_pos + 0.05), Inches(2), Inches(0.3),
                 font_size=11, bold=True, color=COLOR_GRAY)
    add_text_box(slide, font,
                 Inches(6.7), Inches(y_pos + 0.35), Inches(5.8), Inches(0.35),
                 font_size=16, bold=True, color=COLOR_FOREST_GREEN)
    add_text_box(slide, f"來源：{source}  ·  用途：{usage}",
                 Inches(6.7), Inches(y_pos + 0.72), Inches(6.0), Inches(0.3),
                 font_size=10, color=COLOR_GRAY)
    y_pos += 1.25

# 圓角 & 陰影規範
add_rect(slide, Inches(6.5), Inches(5.2), Inches(6.5), Inches(1.9), COLOR_WHITE,
         line_color=COLOR_LIGHT_GREEN, line_width=Pt(1))
add_text_box(slide, "📐  視覺規範",
             Inches(6.7), Inches(5.3), Inches(6.0), Inches(0.35),
             font_size=14, bold=True, color=COLOR_FOREST_GREEN)
add_bullet_text(slide, [
    "圓角：rounded-lg（8px）· 卡片 / 按鈕統一",
    "陰影：shadow-lg · hover:shadow-xl（懸停加深）",
    "間距：優雅留白，py-12 / gap-8 為主要節奏",
    "過渡：transition-colors duration-200（所有互動）",
], Inches(6.7), Inches(5.65), Inches(6.0), Inches(1.3), font_size=11)


# ════════════════════════════════════════════════════════
# 投影片 5：RWD 導覽列架構
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, COLOR_BEIGE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), COLOR_FOREST_GREEN)
add_text_box(slide, "04  RWD 導覽列架構（Golden Navigation Bar）",
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=COLOR_WHITE)

# 桌面版示意
add_text_box(slide, "🖥️  桌面版（md 以上）",
             Inches(0.4), Inches(1.2), Inches(6), Inches(0.4),
             font_size=15, bold=True, color=COLOR_FOREST_GREEN)
add_rect(slide, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.7), COLOR_FOREST_GREEN)
add_text_box(slide, "🌲 小棲蘭",
             Inches(0.5), Inches(1.7), Inches(2), Inches(0.6),
             font_size=14, bold=True, color=COLOR_WHITE)
add_text_box(slide, "最新活動*  |  森林農產  |  日誌  |  🛒(badge)  |  FB",
             Inches(5.5), Inches(1.7), Inches(7.2), Inches(0.6),
             font_size=13, color=COLOR_WHITE, align=PP_ALIGN.RIGHT)
add_text_box(slide, "* Active 狀態：font-bold + border-b-2 border-emerald-900",
             Inches(0.4), Inches(2.4), Inches(12.5), Inches(0.3),
             font_size=10, color=COLOR_GRAY)

# 手機版示意
add_text_box(slide, "📱  手機版（md 以下）",
             Inches(0.4), Inches(2.85), Inches(6), Inches(0.4),
             font_size=15, bold=True, color=COLOR_FOREST_GREEN)
add_rect(slide, Inches(0.4), Inches(3.3), Inches(12.5), Inches(0.7), COLOR_FOREST_GREEN)
add_text_box(slide, "🌲 小棲蘭",
             Inches(0.5), Inches(3.35), Inches(2), Inches(0.6),
             font_size=14, bold=True, color=COLOR_WHITE)
add_text_box(slide, "🛒(badge)  |  FB  |  ☰",
             Inches(9.5), Inches(3.35), Inches(3.2), Inches(0.6),
             font_size=13, color=COLOR_WHITE, align=PP_ALIGN.RIGHT)

# 下拉選單
add_rect(slide, Inches(0.4), Inches(4.0), Inches(12.5), Inches(1.0), COLOR_LIGHT_GREEN)
add_text_box(slide, "📅 最新活動  |  🌿 森林農產  |  📖 日誌",
             Inches(0.6), Inches(4.1), Inches(12), Inches(0.8),
             font_size=13, color=COLOR_FOREST_GREEN)
add_text_box(slide, "↑ 漢堡選單展開後的下拉區域（hidden → 點擊 ☰ 切換）",
             Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.3),
             font_size=10, color=COLOR_GRAY)

# Tailwind 類別表
add_text_box(slide, "🔧  關鍵 Tailwind 類別",
             Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.4),
             font_size=15, bold=True, color=COLOR_FOREST_GREEN)

tw_classes = [
    ("<header>", "sticky top-0 z-50 bg-stone-50 shadow-md"),
    ("桌面連結群組", "hidden md:flex items-center gap-6"),
    ("手機圖示群組", "flex md:hidden items-center gap-4"),
    ("Active 連結", "text-emerald-900 font-bold border-b-2 border-emerald-900 pb-0.5"),
    ("購物車 Badge", "absolute -top-2 -right-2 bg-red-600 text-white text-xs rounded-full w-5 h-5"),
    ("下拉選單", "hidden md:hidden bg-stone-50 border-t border-stone-200 px-4 py-3 flex flex-col"),
]

for i, (element, classes) in enumerate(tw_classes):
    col = i % 2
    row = i // 2
    x = Inches(0.4 + col * 6.5)
    y = Inches(5.85 + row * 0.5)
    add_rect(slide, x, y, Inches(1.8), Inches(0.4), COLOR_FOREST_GREEN)
    add_text_box(slide, element, x, y, Inches(1.8), Inches(0.4),
                 font_size=10, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(2.2 + col * 6.5), y, Inches(4.5), Inches(0.4), COLOR_WHITE,
             line_color=COLOR_LIGHT_GREEN, line_width=Pt(0.5))
    add_text_box(slide, classes,
                 Inches(2.3 + col * 6.5), y, Inches(4.3), Inches(0.4),
                 font_size=9, color=COLOR_CHARCOAL)


# ════════════════════════════════════════════════════════
# 投影片 6：Carousel & 購物車 Badge JS 架構
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, COLOR_BEIGE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), COLOR_FOREST_GREEN)
add_text_box(slide, "05  Carousel & 購物車 Badge — Vanilla ES6 JS 架構",
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=COLOR_WHITE)

# 左欄：Carousel
add_rect(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(6.0), COLOR_WHITE,
         line_color=COLOR_MID_GREEN, line_width=Pt(1.5))
add_text_box(slide, "🎠  Carousel 輪播控制器（events.html）",
             Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=COLOR_FOREST_GREEN)

carousel_code = [
    "// HTML 結構：CSS transition-opacity 驅動",
    ".carousel-slide { transition: opacity 0.7s ease-in-out; }",
    "",
    "// ES6 核心邏輯",
    "const slides = document.querySelectorAll('.carousel-slide');",
    "const dots   = document.querySelectorAll('.carousel-dot');",
    "let current = 0;",
    "",
    "function goToSlide(index) {",
    "  slides[current].classList.replace('opacity-100','opacity-0');",
    "  current = (index + slides.length) % slides.length;",
    "  slides[current].classList.replace('opacity-0','opacity-100');",
    "}",
    "",
    "// 自動輪播（4 秒）",
    "const timer = setInterval(() => goToSlide(current+1), 4000);",
    "",
    "// 解耦設計：移除此 <script> 即可換 WP Slider 插件",
]

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = False
for i, line in enumerate(carousel_code):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.size = Pt(9.5)
    run.font.name = "Consolas"
    if line.startswith("//"):
        run.font.color.rgb = COLOR_MID_GREEN
    elif line.startswith(".") or line.startswith("const") or line.startswith("let") or line.startswith("function"):
        run.font.color.rgb = COLOR_FOREST_GREEN
        run.font.bold = True
    else:
        run.font.color.rgb = COLOR_CHARCOAL

# 右欄：購物車 Badge
add_rect(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(6.0), COLOR_WHITE,
         line_color=COLOR_RED, line_width=Pt(1.5))
add_text_box(slide, "🛒  購物車 Badge 計數器（shop.html）",
             Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=COLOR_FOREST_GREEN)

badge_code = [
    "// 全域計數器（WooCommerce 可直接替換）",
    "let cartCount = 0;",
    "",
    "function addToCart(productName, price) {",
    "  cartCount++;",
    "",
    "  // 更新 Badge（桌面 + 手機）",
    "  const badge = document.getElementById('cart-badge');",
    "  badge.textContent = cartCount;",
    "  badge.classList.remove('hidden');",
    "",
    "  // 彈跳動畫回饋",
    "  badge.classList.add('badge-bounce');",
    "  setTimeout(() => badge.classList.remove('badge-bounce'), 400);",
    "",
    "  // 跨頁持久化（sessionStorage）",
    "  sessionStorage.setItem('cartCount', cartCount);",
    "",
    "  // Toast 提示",
    "  showToast('已加入購物車！');",
    "}",
    "",
    "// WooCommerce 遷移：替換為 wc_add_to_cart_params AJAX",
]

txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.8), Inches(5.0))
tf2 = txBox2.text_frame
tf2.word_wrap = False
for i, line in enumerate(badge_code):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.size = Pt(9.5)
    run.font.name = "Consolas"
    if line.startswith("//"):
        run.font.color.rgb = COLOR_MID_GREEN
    elif line.startswith("let") or line.startswith("function") or line.startswith("const"):
        run.font.color.rgb = COLOR_FOREST_GREEN
        run.font.bold = True
    else:
        run.font.color.rgb = COLOR_CHARCOAL


# ════════════════════════════════════════════════════════
# 投影片 7：CWW 遷移路徑對照表
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, COLOR_BEIGE)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), COLOR_FOREST_GREEN)
add_text_box(slide, "06  CWW 遷移路徑對照表（靜態 MVP → WordPress / WooCommerce）",
             Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.8),
             font_size=20, bold=True, color=COLOR_WHITE)

# 表格標題列
headers = ["靜態 MVP 元素", "WP / WooCommerce 對應", "遷移難度", "說明"]
col_widths = [Inches(3.2), Inches(3.8), Inches(1.5), Inches(4.5)]
x_positions = [Inches(0.3), Inches(3.5), Inches(7.3), Inches(8.8)]

for i, (header, width, x) in enumerate(zip(headers, col_widths, x_positions)):
    add_rect(slide, x, Inches(1.2), width, Inches(0.45), COLOR_FOREST_GREEN)
    add_text_box(slide, header, x, Inches(1.22), width, Inches(0.4),
                 font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# 表格資料
rows = [
    ("href=\"events.html\"", "WP 固定連結 /events/", "⭐ 極低", "直接替換 href 即可"),
    ("href=\"shop.html\"", "WooCommerce Shop /shop/", "⭐ 極低", "WC 自動生成商店頁"),
    ("href=\"blog.html\"", "WP Blog Archive /blog/", "⭐ 極低", "WP 預設文章列表頁"),
    ("href=\"cart.html\"", "WooCommerce Cart /cart/", "⭐ 極低", "WC 自動生成購物車頁"),
    ("Carousel JS", "WP Slider 插件（Swiper/Slick）", "⭐⭐ 低", "移除 <script>，安裝插件"),
    ("addToCart() 函式", "WC AJAX add-to-cart hook", "⭐⭐ 低", "替換為 wc_add_to_cart_params"),
    ("Google Forms 連結", "Gravity Forms / WPForms", "⭐⭐ 低", "安裝表單插件，替換 href"),
    ("YouTube iframe", "Gutenberg 影片區塊", "⭐ 極低", "直接貼 YouTube URL 即可"),
    ("Visa/MC disabled", "WooCommerce 藍新金流 Gateway", "⭐⭐⭐ 中", "安裝金流插件並設定 API"),
    ("sessionStorage 購物車", "WooCommerce 原生購物車", "⭐⭐ 低", "WC 自動接管購物車狀態"),
]

difficulty_colors = {
    "⭐ 極低": RGBColor(0x05, 0x96, 0x69),
    "⭐⭐ 低": COLOR_AMBER,
    "⭐⭐⭐ 中": COLOR_RED,
}

for row_i, (static, wp, difficulty, note) in enumerate(rows):
    y = Inches(1.65 + row_i * 0.5)
    bg = COLOR_WHITE if row_i % 2 == 0 else RGBColor(0xF0, 0xFD, 0xF4)

    for col_i, (text, width, x) in enumerate(zip([static, wp, difficulty, note], col_widths, x_positions)):
        add_rect(slide, x, y, width, Inches(0.48), bg,
                 line_color=RGBColor(0xE5, 0xE7, 0xEB), line_width=Pt(0.5))
        text_color = difficulty_colors.get(text, COLOR_CHARCOAL) if col_i == 2 else COLOR_CHARCOAL
        bold = col_i == 0
        add_text_box(slide, text, Inches(x.inches + 0.05), y, Inches(width.inches - 0.1), Inches(0.48),
                     font_size=10, bold=bold, color=text_color, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════
# 投影片 8：結語 & 下一步行動
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, COLOR_DARK_BG)

# 左側色塊
add_rect(slide, Inches(0), Inches(0), Inches(4.5), SLIDE_H, COLOR_FOREST_GREEN)

add_text_box(slide, "🌲",
             Inches(0.5), Inches(1.0), Inches(3.5), Inches(1.5),
             font_size=72, align=PP_ALIGN.CENTER)

add_text_box(slide, "小棲蘭森林農場",
             Inches(0.3), Inches(2.8), Inches(3.9), Inches(0.6),
             font_size=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "V2 靜態 MVP 完成",
             Inches(0.3), Inches(3.4), Inches(3.9), Inches(0.5),
             font_size=14, color=COLOR_LIGHT_GREEN, align=PP_ALIGN.CENTER)

add_text_box(slide, "5 頁面 · 100% Mobile-First\nCWW 架構對齊 · 解耦設計",
             Inches(0.3), Inches(4.0), Inches(3.9), Inches(0.8),
             font_size=12, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)

# 右側：下一步
add_text_box(slide, "✅  MVP 完成清單",
             Inches(5.0), Inches(1.0), Inches(8.0), Inches(0.5),
             font_size=18, bold=True, color=COLOR_WHITE)

completed = [
    "index.html — 全螢幕歡迎頁（品牌簽名 + 點擊跳轉）",
    "events.html — 活動頁（導覽列 + Carousel + 活動卡片）",
    "shop.html — 電商頁（3 商品 + Badge 計數器 + 支付 Hook）",
    "blog.html — 日誌頁（防爆 YouTube + 情境式 CTA）",
    "cart.html — 購物車（數量控制 + 信任徽章 + 結帳 Modal）",
]
add_bullet_text(slide, completed,
                Inches(5.0), Inches(1.6), Inches(8.0), Inches(2.5),
                font_size=13, color=COLOR_LIGHT_GREEN, bullet="✓ ")

add_text_box(slide, "🚀  下一步行動建議",
             Inches(5.0), Inches(4.2), Inches(8.0), Inches(0.5),
             font_size=18, bold=True, color=COLOR_WHITE)

next_steps = [
    "替換所有 placehold.co 佔位圖為農場實際照片",
    "更新 Google Forms 連結為實際報名表單 URL",
    "更新 Facebook 連結為農場粉絲專頁",
    "填入農場實際電話、地址、聯絡資訊",
    "確認後部署至 Cloudways，開始 WP 遷移規劃",
]
add_bullet_text(slide, next_steps,
                Inches(5.0), Inches(4.8), Inches(8.0), Inches(2.5),
                font_size=13, color=COLOR_AMBER, bullet="→ ")


# ════════════════════════════════════════════════════════
# 儲存 PPTX
# ════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "小棲蘭網站架構說明書.pptx")
prs.save(output_path)
print(f"✅ PPTX 已成功生成：{output_path}")
